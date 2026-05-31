"""build_pptx - Presentacion ejecutiva del CCV.

Diseno LIGHT OPTA matching outputs/viz + PDF del TFM:
- Fondo blanco limpio
- Paleta azul + negro + blanco (matching logo JO + pipeline viz)
- Font Chakra Petch
- Logo JO grande TOP-RIGHT (firma personal)
- Logo SDC pequeno BOTTOM-LEFT (attribution master)
- Imagenes con constraints w/h para no salirse del slide
- Bullets * (no em-dash)
- Pensado para LECTURA, no presentacion oral
"""
from __future__ import annotations
from pathlib import Path

from PIL import Image as _PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = Path(__file__).resolve().parents[1]
OUT  = REPO / 'TFM' / 'CCV_Presentacion_Ejecutiva.pptx'
VIZ  = REPO / 'outputs' / 'viz'
LOGO = REPO / 'outputs' / 'assets' / 'logo.png'
SDC  = REPO / 'outputs' / 'assets' / 'sdc_logo.png'

# Paleta: azul + negro + blanco (matches logo JO + pipeline viz)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
NAVY    = RGBColor(0x1E, 0x3A, 0x8A)  # azul profundo titles + acentos
INK     = RGBColor(0x1A, 0x1A, 0x1A)  # negro suave body
GRAY    = RGBColor(0x55, 0x5B, 0x6E)  # gris caption
SOFT    = RGBColor(0x9C, 0xA3, 0xAF)  # gris very soft
BLUE_LT = RGBColor(0x3B, 0x82, 0xF6)  # azul highlights / cifras hero
SEP     = RGBColor(0xE5, 0xE7, 0xEB)  # gris linea separadora
FONT    = 'Chakra Petch'
BULLET  = '•'  # • bullet

# Zonas reservadas por slide (logo JO top-right, SDC bottom-left)
# Slide: 13.333 x 7.500 inches
LOGO_JO_W = 1.65   # ancho logo JO top-right
LOGO_JO_X = 11.50  # x posicion logo JO (13.333 - 1.65 - 0.18 margen)
LOGO_JO_Y = 0.25
LOGO_SDC_W = 0.45
LOGO_SDC_X = 0.45
LOGO_SDC_Y = 6.85
# Area de contenido segura:
CONTENT_X_MIN = 0.55
CONTENT_X_MAX = 13.10  # parar antes del borde
CONTENT_Y_TOP = 1.90   # debajo de title + separador
CONTENT_Y_BOT = 6.65   # encima del footer (logo SDC + page num)
CONTENT_W = CONTENT_X_MAX - CONTENT_X_MIN  # 12.55
CONTENT_H = CONTENT_Y_BOT - CONTENT_Y_TOP  # 4.75


def _white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _txt(slide, x, y, w, h, text, size, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def _img(slide, path, x, y, w=None, h=None):
    kw = {}
    if w is not None: kw['width']  = Inches(w)
    if h is not None: kw['height'] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def _img_fit(slide, path, x_center, y_top, max_w, max_h):
    """Inserta imagen respetando aspect ratio dentro de max_w x max_h.
    Centra horizontalmente alrededor de x_center. Garantiza que la
    imagen nunca se salga del area dada.
    """
    im = _PILImage.open(path)
    iw, ih = im.size
    ratio = iw / ih
    fit_h = max_w / ratio
    if fit_h <= max_h:
        w = max_w; h = fit_h
    else:
        h = max_h; w = max_h * ratio
    x = x_center - w / 2
    return _img(slide, path, x=x, y=y_top, w=w, h=h)


def _line(slide, x, y, w, color=NAVY, weight_pt=2):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Pt(weight_pt))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


def _chrome(slide, n, total, section=None):
    """Cromados estandar: logo JO top-right + chip top-left + SDC bottom-left + page bottom-right."""
    # logo JO top-right (firma personal)
    _img(slide, LOGO, x=LOGO_JO_X, y=LOGO_JO_Y, w=LOGO_JO_W)
    # chip seccion top-left
    if section:
        _txt(slide, x=0.55, y=0.40, w=4.0, h=0.30,
             text=section.upper(), size=10, color=NAVY, bold=True,
             align=PP_ALIGN.LEFT)
    # logo SDC bottom-left
    _img(slide, SDC, x=LOGO_SDC_X, y=LOGO_SDC_Y, w=LOGO_SDC_W)
    # page num bottom-right
    _txt(slide, x=12.20, y=7.10, w=0.85, h=0.30,
         text=f'{n:02d} / {total:02d}', size=10, color=SOFT,
         align=PP_ALIGN.RIGHT)


def _title(slide, text, size=32):
    # title pegado al borde izquierdo, ancho hasta antes del logo
    _txt(slide, x=0.55, y=0.95, w=10.90, h=0.85,
         text=text, size=size, color=NAVY, bold=True)
    _line(slide, x=0.55, y=1.75, w=0.55, color=NAVY, weight_pt=3)


def _bullet(slide, x, y, w, h, items, size=15, color=INK, indent=0.18):
    """Lista de bullets con marker • en color navy + texto en INK."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.20
        p.space_before = Pt(0)
        p.space_after = Pt(6)
        # marker
        rm = p.add_run()
        rm.text = BULLET + '  '
        rm.font.name = FONT
        rm.font.size = Pt(size)
        rm.font.bold = True
        rm.font.color.rgb = NAVY
        # text
        rt = p.add_run()
        rt.text = line
        rt.font.name = FONT
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return tb


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.500)

    total = 27

    # ============ 1. PORTADA ============
    s = _blank(prs)
    # logo SDC arriba (en portada va mas grande)
    _img(s, SDC, x=6.34, y=0.40, w=0.65)
    _line(s, x=4.5, y=1.65, w=4.3, color=NAVY, weight_pt=4)
    _txt(s, 0.5, 1.95, 12.3, 1.0, 'CAUSAL CLUTCH VALUE',
         size=58, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 0.5, 3.05, 12.3, 0.6,
         'Cuantificación causal del clutch en el fútbol de élite',
         size=22, color=BLUE_LT, align=PP_ALIGN.CENTER)
    _txt(s, 0.5, 4.20, 12.3, 0.5,
         'Trabajo Fin de Máster',
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, 0.5, 4.70, 12.3, 0.4,
         'Máster en Big Data Aplicado al Scouting Deportivo',
         size=13, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, 0.5, 5.10, 12.3, 0.4,
         'Sports Data Campus · 2025-2026',
         size=13, color=GRAY, align=PP_ALIGN.CENTER)
    _img(s, LOGO, x=5.42, y=5.85, w=2.50)

    # ============ 2. RESUMEN EJECUTIVO ============
    s = _blank(prs)
    _title(s, 'En una página')
    _txt(s, 0.55, 2.00, 12.5, 0.50,
         'Pregunta', size=16, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 2.50, 12.5, 0.55,
         '¿Puede medirse causalmente la respuesta individual del jugador al shock emocional del gol,',
         size=17, color=INK)
    _txt(s, 0.55, 2.95, 12.5, 0.55,
         'separada del empuje colectivo de su propio equipo?',
         size=17, color=INK)
    _txt(s, 0.55, 3.75, 12.5, 0.50,
         'Marco', size=16, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 4.25, 12.5, 0.55,
         'CCV: pipeline de 6 fases con identificación causal en 5 capas',
         size=17, color=INK)
    _txt(s, 0.55, 4.70, 12.5, 0.55,
         '(DiD within-player + AIPW cross-fitting + CATE jerárquico bayesiano multivariate).',
         size=17, color=INK)
    _txt(s, 0.55, 5.50, 12.5, 0.50,
         'Hallazgo', size=16, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 6.00, 12.5, 0.55,
         'ATE poblacional ≈ 0, pero el efecto vive en la heterogeneidad individual:',
         size=17, color=INK)
    _txt(s, 0.55, 6.45, 12.5, 0.55,
         '4 pressure-clutch leaders + 22 jugadores con respuesta ofensiva significativa post-marcar.',
         size=17, color=INK)
    _chrome(s, 2, total, 'Resumen')

    # ============ 3. MOTIVACION ============
    s = _blank(prs)
    _title(s, 'Por qué este TFM')
    _txt(s, 0.55, 2.00, 12.5, 0.55,
         'El clutch existe en la intuición del aficionado y del scout.',
         size=20, color=INK)
    _txt(s, 0.55, 2.55, 12.5, 0.55,
         'Mbappé marca tres goles en la Final. Lautaro falla un mano-a-mano. Saka se ofrece al penalti.',
         size=16, color=GRAY)
    _txt(s, 0.55, 3.30, 12.5, 0.55,
         'Pero los rankings públicos miden lo que pasa, no la reacción al shock.',
         size=20, color=INK)
    _txt(s, 0.55, 3.85, 12.5, 0.55,
         'Parten de eventing puro, sin tracking sincronizado y sin separar lo individual del bloque.',
         size=16, color=GRAY)
    _txt(s, 0.55, 4.65, 12.5, 0.55,
         'Resultado: un jugador con buena estadística post-gol suele estar',
         size=20, color=INK)
    _txt(s, 0.55, 5.20, 12.5, 0.55,
         'en un equipo que ya estaba dominando el partido.',
         size=20, color=INK)
    _line(s, x=0.55, y=6.10, w=4.0, color=BLUE_LT, weight_pt=2)
    _txt(s, 0.55, 6.25, 12.5, 0.55,
         'El clutch verdadero requiere descontar lo que hace su bloque.',
         size=20, color=NAVY, bold=True)
    _chrome(s, 3, total, 'Motivación')

    # ============ 4. PREGUNTA + OBJETIVO ============
    s = _blank(prs)
    _title(s, 'Pregunta de investigación y objetivo')
    _txt(s, 0.55, 2.00, 12.5, 0.50, 'Pregunta',
         size=16, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 2.55, 12.5, 0.55,
         '¿Es la respuesta del jugador al shock emocional del gol',
         size=21, color=NAVY, bold=True)
    _txt(s, 0.55, 3.10, 12.5, 0.55,
         'una propiedad individual estable, residualizada contra el bloque?',
         size=21, color=NAVY, bold=True)
    _txt(s, 0.55, 4.20, 12.5, 0.50, 'Objetivo general',
         size=16, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 4.75, 12.5, 0.55,
         'Estimar causalmente el efecto del shock sobre el rendimiento individual,',
         size=19, color=INK)
    _txt(s, 0.55, 5.25, 12.5, 0.55,
         'sobre cuatro canales (Ofensivo · Defensivo · Off-ball · Físico)',
         size=19, color=INK)
    _txt(s, 0.55, 5.75, 12.5, 0.55,
         'y tres contextos (post-gol a favor · post-gol en contra · proximidad de eliminación),',
         size=19, color=INK)
    _txt(s, 0.55, 6.25, 12.5, 0.55,
         'con cuantificación de la incertidumbre por jugador.',
         size=19, color=INK)
    _chrome(s, 4, total, 'Objetivo')

    # ============ 5. 6 OBJETIVOS ESPECIFICOS ============
    s = _blank(prs)
    _title(s, 'Seis objetivos específicos')
    oes = [
        ('OE1', 'Backbone temporal: Win Probability bayesiana + leverage + proximidad de eliminación'),
        ('OE2', 'Post-Shot xG calibrado + corpus de near-miss como control cuasi-experimental'),
        ('OE3', 'Cuantificar el rendimiento en cuatro canales independientes'),
        ('OE4', 'Identificación causal dual: DiD within-player + AIPW con cross-fitting'),
        ('OE5', 'CATE jerárquico bayesiano multivariate (NUTS HMC + LKJ cross-canal)'),
        ('OE6', 'Ensamblar tabla scout-facing por jugador, canal y contexto'),
    ]
    for i, (k, v) in enumerate(oes):
        y = 2.10 + i * 0.72
        _txt(s, 0.65, y, 1.0, 0.55, k, size=20, color=BLUE_LT, bold=True)
        _txt(s, 1.85, y+0.05, 10.7, 0.55, v, size=16, color=INK)
    _chrome(s, 5, total, 'Objetivo')

    # ============ 6. TRES FUENTES ============
    s = _blank(prs)
    _title(s, 'Tres fuentes de datos')
    cols = [
        ('PFF FC',     'World Cup Qatar 2022',  'TARGET',
         '64 partidos · eventing + tracking 25 Hz sincronizado · grades humanos por evento · rosters oficiales'),
        ('Wyscout',    '2017/18 open dataset',  'WP BACKBONE',
         '1941 partidos · cinco grandes ligas europeas + WC 2018 + Euro 2016'),
        ('StatsBomb',  'Open data filtrado',    'PSxG + VAEP',
         '136 partidos · Euro 2020 + Euro 2024 + Bundesliga 23/24 · incluye 360 freeze-frames'),
    ]
    for i, (k, sub, role, vol) in enumerate(cols):
        y = 2.05 + i * 1.55
        _txt(s, 0.65, y, 3.0, 0.55, k, size=24, color=NAVY, bold=True)
        _txt(s, 3.65, y+0.05, 4.5, 0.55, sub, size=14, color=GRAY, italic=True)
        _txt(s, 9.30, y+0.05, 3.4, 0.55, role, size=13, color=BLUE_LT, bold=True, align=PP_ALIGN.RIGHT)
        _txt(s, 0.65, y+0.65, 12.0, 0.75, vol, size=13, color=INK)
        if i < 2:
            _line(s, x=0.65, y=y+1.40, w=11.9, color=SEP, weight_pt=1)
    _chrome(s, 6, total, 'Datos')

    # ============ 7. POR QUE PFF ============
    s = _blank(prs)
    _title(s, 'Por qué PFF FC')
    _txt(s, 0.55, 2.00, 12.5, 0.55,
         'Único corpus público que combina los tres ingredientes estructurales',
         size=19, color=INK)
    _txt(s, 0.55, 2.55, 12.5, 0.55,
         'que la pregunta de investigación exige simultáneamente:',
         size=19, color=INK)
    ings = [
        ('1', 'Tracking continuo full-pitch', 'de los 22 jugadores a 25 Hz'),
        ('2', 'Sincronización exacta',         'tracking ↔ eventing vía game_event.start_frame'),
        ('3', 'Cobertura completa',            'de fase eliminatoria con prórroga incluida'),
    ]
    for i, (n, h, sub) in enumerate(ings):
        y = 3.45 + i * 0.85
        _txt(s, 0.85, y, 0.8, 0.7, n, size=40, color=BLUE_LT, bold=True)
        _txt(s, 1.90, y+0.05, 4.5, 0.5, h, size=19, color=NAVY, bold=True)
        _txt(s, 6.45, y+0.10, 6.5, 0.5, sub, size=15, color=INK)
    _txt(s, 0.55, 6.20, 12.5, 0.50,
         'Los grades humanos por evento son ingrediente opcional (prior CATE + validación externa), no estructural.',
         size=12, color=GRAY, italic=True)
    _chrome(s, 7, total, 'Datos')

    # ============ 8. VOLUMEN ============
    s = _blank(prs)
    _title(s, 'Corpus de trabajo en cifras')
    cells = [
        ('64',  'partidos',     'del Mundial 2022'),
        ('172', 'shocks',       '120 grupos + 52 eliminatoria'),
        ('70',  'near-misses',  '5 tipologías pre-registradas'),
        ('511', 'jugadores',    'con minutos en campo'),
        ('233', 'jugadores',    'con ≥ 270 min (muestra fiable)'),
        ('166', 'jugadores',    'con grades PFF para validación'),
    ]
    for i, (n, label, sub) in enumerate(cells):
        row, col = divmod(i, 3)
        x = 0.50 + col * 4.25
        y = 2.05 + row * 2.30
        _txt(s, x, y,      4.20, 1.10, n, size=68, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, x, y+1.20, 4.20, 0.45, label, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, x, y+1.65, 4.20, 0.45, sub, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    _chrome(s, 8, total, 'Datos')

    # ============ 9. PIPELINE DAG ============
    s = _blank(prs)
    _title(s, 'Pipeline en seis fases')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'Extracción → WP backbone → shocks/near-miss → cuatro canales en paralelo → CATE jerárquico → ensamblaje scout-facing',
         size=11, color=GRAY, italic=True)
    _img_fit(s, VIZ / 'fig_cap4_pipeline_dag.jpg',
             x_center=6.665, y_top=2.30, max_w=CONTENT_W, max_h=4.30)
    _chrome(s, 9, total, 'Pipeline')

    # ============ 10. MAPA CONCEPTUAL ============
    s = _blank(prs)
    _title(s, 'Mapa conceptual del CCV')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'Tipo de shock × ¿el jugador rompe del bloque? × ¿reacción positiva? → proyectado sobre cuatro canales',
         size=11, color=GRAY, italic=True)
    _img_fit(s, VIZ / 'fig_cap4_mapa_conceptual.png',
             x_center=6.665, y_top=2.30, max_w=CONTENT_W, max_h=4.30)
    _chrome(s, 10, total, 'Marco')

    # ============ 11. 5 CAPAS CAUSALES ============
    s = _blank(prs)
    _title(s, 'Cinco capas de identificación causal')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'Cada capa retira un confusor distinto. El estimador causal emerge sólo cuando las cinco operan en secuencia.',
         size=11, color=GRAY, italic=True)
    _img_fit(s, VIZ / 'fig_cap4_capas_causales.png',
             x_center=6.665, y_top=2.30, max_w=CONTENT_W, max_h=4.30)
    _chrome(s, 11, total, 'Marco')

    # ============ 12. DiD ============
    s = _blank(prs)
    _title(s, 'Diferencias-en-diferencias within-player')
    _txt(s, 0.55, 1.95, 5.8, 0.45, 'Modelo TWFE',
         size=17, color=BLUE_LT, bold=True)
    _txt(s, 0.55, 2.45, 5.8, 0.55,
         'Y = αᵢ,ₛ + γₜ−ₜₛ + τ · 1[t > tₛ] + ε',
         size=16, color=INK, italic=True)
    _txt(s, 0.55, 3.30, 5.8, 0.45, 'Verificaciones independientes',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 0.55, 3.85, 5.9, 3.0, [
        'Sun & Abraham 2021 — sin contaminación cohortes',
        'Borusyak-Jaravel-Spiess 2024 — imputación contrafactual',
        'HonestDiD Rambachan-Roth 2023 — cotas de relajación',
        'Test F de pre-tendencias agregadas (Roth 2022)',
    ], size=13)
    _img_fit(s, VIZ / 'event_study.png',
             x_center=10.10, y_top=2.00, max_w=6.30, max_h=4.40)
    _chrome(s, 12, total, 'Identificación causal')

    # ============ 13. AIPW ============
    s = _blank(prs)
    _title(s, 'AIPW con cross-fitting (DML)')
    _txt(s, 0.55, 2.00, 12.5, 0.55,
         'Doblemente robusto asintótico bajo tres condiciones:',
         size=19, color=INK)
    _bullet(s, 0.65, 2.70, 12.4, 1.8, [
        'correcta especificación de al menos uno de los nuisance',
        'overlap positivo (0 < ê(X) < 1 con margen)',
        'consistencia n⁻¹ᐟ⁴ de los nuisance (Chernozhukov 2018)',
    ], size=16)
    _txt(s, 0.55, 4.70, 12.5, 0.50, 'Implementación',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 0.65, 5.25, 12.4, 1.5, [
        'Modelos nuisance: LightGBM, 5-fold estratificado por partido',
        'Propensity score ê(X) + outcome regression μ̂(X) sobre confounders observables',
        'Confounders: minuto, marcador, fuerza del rival, fase del torneo',
    ], size=14)
    _chrome(s, 13, total, 'Identificación causal')

    # ============ 14. CATE BAYESIANO ============
    s = _blank(prs)
    _title(s, 'CATE jerárquico bayesiano multivariate')
    _txt(s, 0.55, 1.95, 5.8, 0.45, 'Estructura',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 0.55, 2.45, 5.9, 2.0, [
        '4 canales conjuntamente (no separados)',
        'Correlaciones cross-canal vía LKJ Cholesky',
        'Priors informados por grades PFF (opcionales)',
        'Parametrización non-centered (Betancourt 2015)',
    ], size=13)
    _txt(s, 0.55, 4.50, 5.8, 0.45, 'Sampling',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 0.55, 5.00, 5.9, 2.0, [
        'NUTS HMC · 4 chains × 1000 warmup + 1000 sampling',
        'target_accept_prob = 0.95 (geometría LKJ)',
        '0 divergencias · 108 / 144 hyperparams R-hat < 1.05',
        'PPC 8/8 canales calibrados (KS p > 0.05)',
    ], size=13)
    _img_fit(s, VIZ / 'fig_cap6_cate_heterogeneity.png',
             x_center=10.10, y_top=2.00, max_w=6.30, max_h=4.40)
    _chrome(s, 14, total, 'CATE')

    # ============ 15-18. 4 CANALES ============
    canales = [
        (15, 'Canal Ofensivo',
            [('atomic-VAEP (CatBoost) sobre representación SPADL', 'AUC scores 0.83 / concedes 0.87'),
             ('un-xPass (Robberechts 2023) — creatividad en pase', 'AUC 0.83'),
             ('Entrenamiento: 456k acciones + 137k pases',         'StatsBomb open data')]),
        (16, 'Canal Defensivo',
            [('VDEP estricto (Toda 2022): recovery + attacked',    'AUC 0.80 / 0.83'),
             ('exPress (Lee 2025) — P(recuperación < 5s | press)', 'AUC 0.62'),
             ('Atribución frame-nearest al defensor más cercano',  'Frame level 25 Hz')]),
        (17, 'Canal Off-ball',
            [('OBSO + C-OBSO (Spearman 2018, Teranishi 2022)',     'Off-ball Score Opportunity'),
             ('Pitch Control vectorizado sobre tracking 25 Hz',     'C-OBSO vs grade PFF r = +0.29'),
             ('64 partidos full-time 22 jugadores · 25 Hz',         'n=673, p < 10⁻¹⁴')]),
        (18, 'Canal Físico',
            [('Protocolo Bradley 2024 sobre velocidades suavizadas','Filtro Butterworth fase cero'),
             ('Cap a 11 m/s preservando dirección',                 '64 partidos a 25 Hz'),
             ('Residualización bayesiana jerárquica via SVI',        'Línea base personal por curva minuto')]),
    ]
    for n, title, rows in canales:
        s = _blank(prs)
        _title(s, title)
        _txt(s, 0.65, 2.00, 8.5, 0.40, 'Building block', size=12, color=GRAY, bold=True)
        _txt(s, 9.30, 2.00, 3.6, 0.40, 'Métrica / Validación', size=12, color=GRAY, bold=True)
        _line(s, x=0.65, y=2.40, w=12.0, color=SEP, weight_pt=1)
        for i, (bb, met) in enumerate(rows):
            y = 2.55 + i * 1.15
            _txt(s, 0.65, y, 8.5, 0.55, bb, size=16, color=INK)
            _txt(s, 9.30, y, 3.6, 0.55, met, size=14, color=BLUE_LT, bold=True)
            if i < len(rows) - 1:
                _line(s, x=0.65, y=y+1.00, w=12.0, color=SEP, weight_pt=1)
        _chrome(s, n, total, 'Canales')

    # ============ 19. PSxG ANCLA ============
    s = _blank(prs)
    _title(s, 'PSxG: ancla técnica del corpus near-miss')
    _txt(s, 0.55, 1.95, 5.8, 0.45,
         'Modelo predictivo',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 0.55, 2.45, 5.9, 1.8, [
        'LightGBM + Optuna (60 iter) + isotonic',
        'Features: end_y, end_z, configuración 360 freeze-frame',
        'Train: StatsBomb open data fuera del Mundial',
    ], size=13)
    _txt(s, 0.55, 4.30, 5.8, 0.45, 'Métricas',
         size=17, color=BLUE_LT, bold=True)
    metrics = [
        ('AUC OOF',                       '0.968'),
        ('AUC holdout WC22',              '0.976'),
        ('vs StatsBomb xG (pre-disparo)', '0.844'),
        ('ECE holdout',                   '0.011'),
        ('Brier holdout',                 '0.037'),
    ]
    for i, (k, v) in enumerate(metrics):
        y = 4.85 + i * 0.36
        _txt(s, 0.65, y, 4.5, 0.35, k, size=13, color=INK)
        _txt(s, 5.15, y, 1.4, 0.35, v, size=14, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    _img_fit(s, VIZ / 'fig_cap6_psxg_calibration.png',
             x_center=10.10, y_top=2.00, max_w=6.30, max_h=4.40)
    _chrome(s, 19, total, 'Validación')

    # ============ 20. FICHA MBAPPE ============
    s = _blank(prs)
    _title(s, 'Ficha CCV — Kylian Mbappé')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'Pressure-clutch leader del torneo · CATE +0.110 desv. estándar · P(β > 0) = 0.97',
         size=13, color=BLUE_LT, bold=True)
    _img_fit(s, VIZ / 'radar_report_3870.png',
             x_center=6.665, y_top=2.35, max_w=CONTENT_W, max_h=4.25)
    _chrome(s, 20, total, 'Output scout-facing')

    # ============ 21. SCATTER REMONTADOR x CERROJO ============
    s = _blank(prs)
    _title(s, 'Remontador × Cerrojo — los 511 jugadores')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'Dos índices ortogonales: respuesta al gol en contra (Remontador) vs. respuesta al gol a favor (Cerrojo).',
         size=11, color=GRAY, italic=True)
    _img_fit(s, VIZ / 'scatter_remontador_cerrojo.png',
             x_center=6.665, y_top=2.30, max_w=CONTENT_W, max_h=4.30)
    _chrome(s, 21, total, 'Output scout-facing')

    # ============ 22. ATE = 0 ============
    s = _blank(prs)
    _title(s, 'ATE poblacional indistinguible de cero')
    _img_fit(s, VIZ / 'fig_cap6_window_sensitivity.png',
             x_center=4.20, y_top=2.00, max_w=7.50, max_h=4.50)
    _txt(s, 8.40, 2.00, 4.4, 0.45, 'Lectura',
         size=17, color=BLUE_LT, bold=True)
    _bullet(s, 8.40, 2.55, 4.5, 1.5, [
        '8 celdas canal × contexto',
        'Ninguna cruza significatividad al 5%',
    ], size=13)
    _txt(s, 8.40, 4.00, 4.4, 0.45, 'Robustez', size=17, color=BLUE_LT, bold=True)
    _bullet(s, 8.40, 4.55, 4.5, 1.5, [
        'Placebo 1000 perm + BH FDR: ningún canal sale del null (p_FDR = 0.98)',
        'HonestDiD: cotas contienen el cero en los 3 niveles',
    ], size=11)
    _line(s, x=8.40, y=6.15, w=4.0, color=NAVY, weight_pt=2)
    _txt(s, 8.40, 6.25, 4.5, 0.55,
         'El efecto medio no se',
         size=14, color=NAVY, bold=True)
    _txt(s, 8.40, 6.55, 4.5, 0.55,
         'distingue del ruido.',
         size=14, color=NAVY, bold=True)
    _chrome(s, 22, total, 'Resultados')

    # ============ 23. 4 LEADERS ============
    s = _blank(prs)
    _title(s, 'Cuatro pressure-clutch leaders identificados')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'CATE pressure-clutch · |P(β)| ≥ 0.85 · sobre 511 jugadores evaluados',
         size=11, color=GRAY, italic=True)
    headers = ['Jugador', 'Selección', 'CATE', 'P(β > 0)']
    rows = [
        ('Kylian Mbappé',    'Francia',    '+0.110', '0.97'),
        ('Marcus Rashford',  'Inglaterra', '+0.050', '0.89'),
        ('Bukayo Saka',      'Inglaterra', '+0.050', '0.86'),
        ('Mohammed Muntari', 'Catar',      '−0.055', '0.12'),
    ]
    x_cols = [0.70, 4.70, 8.30, 10.70]
    w_cols = [4.00, 3.60, 2.40, 2.00]
    y_h = 2.40
    for i, h in enumerate(headers):
        _txt(s, x_cols[i], y_h, w_cols[i], 0.40, h,
             size=13, color=GRAY, bold=True)
    _line(s, x=0.70, y=2.80, w=12.0, color=SEP, weight_pt=1)
    for j, (name, team, cate, prob) in enumerate(rows):
        y = 3.05 + j * 0.85
        _txt(s, x_cols[0], y, w_cols[0], 0.60, name, size=20, color=INK, bold=True)
        _txt(s, x_cols[1], y, w_cols[1], 0.60, team, size=17, color=INK)
        # Muntari: navy mas oscuro para el negativo
        col = BLUE_LT if j < 3 else NAVY
        _txt(s, x_cols[2], y, w_cols[2], 0.60, cate, size=22, color=col, bold=True)
        _txt(s, x_cols[3], y, w_cols[3], 0.60, prob, size=22, color=col, bold=True)
        if j < 3:
            _line(s, x=0.70, y=y+0.75, w=12.0, color=SEP, weight_pt=1)
    _txt(s, 0.55, 6.55, 12.5, 0.40,
         'Muntari es el único negativo significativo: cae bajo presión de eliminación inminente.',
         size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    _chrome(s, 23, total, 'Resultados')

    # ============ 24. 22 OFENSIVO ============
    s = _blank(prs)
    _title(s, 'Ofensivo-tras-marcar: 22 jugadores significativos')
    _txt(s, 0.55, 1.85, 12.5, 0.40,
         'La celda canal × contexto con más señal individual del torneo.',
         size=13, color=GRAY, italic=True)
    _txt(s, 0.55, 2.40, 6.20, 1.4, '5', size=100, color=BLUE_LT, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 4.05, 6.20, 0.55, 'al alza', size=22, color=INK, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 4.65, 6.20, 0.45, 'ascienden producción ofensiva', size=13, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 5.10, 6.20, 0.45, 'tras marcar gol', size=13, color=GRAY, align=PP_ALIGN.CENTER)

    _txt(s, 6.60, 2.40, 6.20, 1.4, '17', size=100, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 6.60, 4.05, 6.20, 0.55, 'a la baja', size=22, color=INK, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 6.60, 4.65, 6.20, 0.45, 'se relajan tras marcar', size=13, color=GRAY, align=PP_ALIGN.CENTER)
    _txt(s, 6.60, 5.10, 6.20, 0.45, '(efecto "ya está hecho")', size=13, color=GRAY, align=PP_ALIGN.CENTER)

    _line(s, x=2.0, y=5.95, w=9.3, color=SEP, weight_pt=1)
    _txt(s, 0.55, 6.10, 12.5, 0.55,
         'Signos opuestos en magnitudes similares → el promedio cancela y el ATE colapsa a cero.',
         size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _chrome(s, 24, total, 'Resultados')

    # ============ 25. CONCLUSION CLAVE ============
    s = _blank(prs)
    _line(s, x=4.5, y=1.65, w=4.3, color=NAVY, weight_pt=4)
    _txt(s, 0.55, 2.15, 12.5, 1.10,
         'ATE nulo no implica efecto nulo.',
         size=44, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 3.50, 12.5, 0.65,
         'El clutch vive en la heterogeneidad individual,',
         size=26, color=BLUE_LT, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 4.15, 12.5, 0.65,
         'no en la media poblacional.',
         size=26, color=BLUE_LT, align=PP_ALIGN.CENTER)
    _line(s, x=5.6, y=5.15, w=2.1, color=SEP, weight_pt=1)
    _txt(s, 0.55, 5.50, 12.5, 0.55,
         'El framework permite identificar QUIÉN responde al shock,',
         size=17, color=INK, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 6.05, 12.5, 0.55,
         'en qué canal, en qué contexto, y con cuánta certeza.',
         size=17, color=INK, align=PP_ALIGN.CENTER)
    _chrome(s, 25, total, 'Conclusión')

    # ============ 26. LIMITACIONES + FUTURO ============
    s = _blank(prs)
    _title(s, 'Limitaciones y líneas futuras')
    _txt(s, 0.65, 2.00, 5.9, 0.45, 'Limitaciones', size=18, color=NAVY, bold=True)
    _bullet(s, 0.65, 2.55, 5.9, 4.0, [
        'Un único torneo (Mundial Qatar 2022)',
        '172 shocks → CATE individual fiable pero no por partido',
        '70 near-misses → triangulación agregada, no por jugador',
        'Building blocks distintos a VAEP son reimplementación propia',
    ], size=14)
    _txt(s, 6.85, 2.00, 5.9, 0.45, 'Líneas futuras', size=18, color=BLUE_LT, bold=True)
    _bullet(s, 6.85, 2.55, 5.9, 4.0, [
        'Más torneos KO (Eurocopa, Champions, World Cup 2026)',
        'Corpus near-miss multitorneo para verificación por jugador',
        'Catálogo de shock ampliado (sustituciones, sup./inf. numérica)',
        'Portabilidad a "clutch del bloque" a nivel equipo',
    ], size=14)
    _chrome(s, 26, total, 'Cierre')

    # ============ 27. GRACIAS ============
    s = _blank(prs)
    _img(s, SDC, x=6.34, y=0.40, w=0.65)
    _line(s, x=4.5, y=1.55, w=4.3, color=NAVY, weight_pt=4)
    _txt(s, 0.55, 1.95, 12.5, 1.20,
         'Gracias',
         size=78, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 3.65, 12.5, 0.55,
         'Documento TFM completo + pipeline E2E reproducible:',
         size=18, color=INK, align=PP_ALIGN.CENTER)
    _txt(s, 0.55, 4.25, 12.5, 0.65,
         'github.com/jaime-oriol/CCV',
         size=26, color=BLUE_LT, bold=True, align=PP_ALIGN.CENTER)
    _img(s, LOGO, x=5.42, y=5.30, w=2.50)
    _txt(s, 0.55, 7.05, 12.5, 0.35,
         'Jaime Oriol Goicoechea · Máster Big Data Aplicado al Scouting Deportivo · Sports Data Campus · 2026',
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

    prs.save(str(OUT))
    print(f'guardado: {OUT.relative_to(REPO)}')
    print(f'  slides: {len(prs.slides)} · size {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f}')


if __name__ == '__main__':
    build()
